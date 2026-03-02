from fastapi import FastAPI, UploadFile, File
from sklearn.feature_extraction.text import TfidfVectorizer
import numpy as np
from pptx import Presentation
from transformers import pipeline
import shutil
import os
import re
from PIL import Image
import pytesseract
import io


app = FastAPI()
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

summarizer = pipeline(
    "summarization",
    model="facebook/bart-large-cnn"
)

def clean_text(text):
    text = re.sub(r'http\S+', '', text)

    text = re.sub(r'\S+@\S+', '', text)

    text = re.sub(r'©.*?reserved\.', '', text, flags=re.IGNORECASE)
    text = re.sub(r'\(c\).*?reserved\.', '', text, flags=re.IGNORECASE)

    text = re.sub(r'.*All rights reserved.*', '', text, flags=re.IGNORECASE)

    text = re.sub(r'©?\s?\d{4}.*', '', text)

    text = re.sub(r'<.*?>', '', text)

    text = re.sub(r'[^A-Za-z0-9.,()\-:\n ]+', ' ', text)

    text = re.sub(r'\s+', ' ', text)

    return text.strip()

def is_reference_slide(text):
    text_lower = text.lower()
    if "references" in text_lower:
        return True
    if text_lower.count("isbn") > 1:
        return True
    return False

def group_short_slides(slides, min_words=60):
    grouped_slides = []
    buffer_text = ""
    buffer_numbers = []

    for slide in slides:
        content = clean_text(slide["content"])
        word_count = len(content.split())

        if word_count == 0:
            continue

        buffer_text += " " + content
        buffer_numbers.append(slide["slide_number"])

        if len(buffer_text.split()) >= min_words:
            grouped_slides.append({
                "slide_numbers": buffer_numbers.copy(),
                "content": buffer_text.strip()
            })
            buffer_text = ""
            buffer_numbers = []

    # Sisa terakhir
    if buffer_text:
        grouped_slides.append({
            "slide_numbers": buffer_numbers,
            "content": buffer_text.strip()
        })

    return grouped_slides

def extract_text_from_ppt(file_path):
    prs = Presentation(file_path)
    slides_text = []

    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_content = []

        for shape in slide.shapes:
            # Ambil text biasa
            if hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text.strip())

            # Ambil gambar dan OCR
            if shape.shape_type == 13:  # 13 = Picture
                image_bytes = shape.image.blob
                try:
                    image = Image.open(io.BytesIO(image_bytes))

                    # Convert ke RGB supaya aman
                    if image.mode != "RGB":
                        image = image.convert("RGB")

                    ocr_text = pytesseract.image_to_string(image)

                    if ocr_text.strip():
                        slide_content.append(ocr_text.strip())

                except Exception as e:
                    print(f"OCR error on slide {slide_number}: {e}")

        full_slide_text = "\n".join(slide_content)

        slides_text.append({
            "slide_number": slide_number,
            "content": full_slide_text
        })

    return slides_text


def summarize_per_slide(slides):
    slide_summaries = []

    grouped = group_short_slides(slides, min_words=80)

    for group in grouped:
        try:
            result = summarizer(
                group["content"],
                max_length=100,
                min_length=40,
                num_beams=6,
                repetition_penalty=2.0,
                early_stopping=True
            )

            summary_text = result[0]["summary_text"]

        except Exception as e:
            summary_text = f"Error summarizing slides {group['slide_numbers']}: {str(e)}"

        slide_summaries.append({
            "slide_numbers": group["slide_numbers"],
            "summary": summary_text
        })

    return slide_summaries


@app.post("/summarize")
async def summarize_ppt(file: UploadFile = File(...)):

    os.makedirs("temp", exist_ok=True)
    file_location = f"temp/{file.filename}"
    with open(file_location, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    slides = extract_text_from_ppt(file_location)

    if not slides:
        return {"message": "No slides found in PPT."}

    slide_summaries = summarize_per_slide(slides)

    return {
        "total_slides": len(slide_summaries),
        "slides_summary": slide_summaries
    }